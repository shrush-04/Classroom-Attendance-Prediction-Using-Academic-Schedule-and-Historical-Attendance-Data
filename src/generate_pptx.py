"""
generate_pptx.py
Creates presentation/project_presentation.pptx from slide data.
Uses python-pptx. All metric values sourced from actual result files.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE_DIR = r"d:\Data_Science_attendence_project"
CHARTS   = os.path.join(BASE_DIR, "outputs", "charts")
OUT_PATH = os.path.join(BASE_DIR, "presentation", "project_presentation.pptx")

# ─── Colour palette ───────────────────────────────────────────────────────────
C_DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy
C_ACCENT    = RGBColor(0x0F, 0x3C, 0x78)   # deep blue
C_HIGHLIGHT = RGBColor(0x00, 0xB4, 0xD8)   # cyan
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xE0, 0xF2, 0xFE)
C_WARN      = RGBColor(0xFF, 0xC3, 0x00)   # amber (for best model highlight)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ─── Helper functions ─────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, l, t, w, h, font_size=18, bold=False,
             color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_bullet_box(slide, lines, l, t, w, h, font_size=16,
                   title_color=C_HIGHLIGHT, body_color=C_WHITE):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.name  = "Calibri"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = title_color if line.startswith("▶") else body_color
    return txb


def bg(slide):
    """Draw dark background."""
    add_rect(slide, 0, 0, 13.33, 7.5, C_DARK_BG)


def header_bar(slide, title_text, subtitle_text=""):
    add_rect(slide, 0, 0, 13.33, 1.35, C_ACCENT)
    add_text(slide, title_text,  0.25, 0.08, 12, 0.75, font_size=30,
             bold=True, color=C_HIGHLIGHT)
    if subtitle_text:
        add_text(slide, subtitle_text, 0.25, 0.85, 12, 0.45, font_size=14,
                 color=C_LIGHT)


def footer(slide, note=""):
    add_rect(slide, 0, 7.05, 13.33, 0.45, RGBColor(0x0F, 0x3C, 0x78))
    msg = ("⚠ SYNTHETIC DATA — No real student names, roll numbers, or email IDs used.  |  " + note
           if note else "⚠ SYNTHETIC DATA — No real student names, roll numbers, or email IDs used.")
    add_text(slide, msg, 0.2, 7.08, 12.9, 0.35, font_size=9,
             color=C_LIGHT, align=PP_ALIGN.CENTER)


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def try_add_image(slide, fname, l, t, w, h):
    path = os.path.join(CHARTS, fname)
    if os.path.isfile(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        add_rect(slide, l, t, w, h, RGBColor(0x22, 0x22, 0x44))
        add_text(slide, f"[Chart: {fname}]", l+0.1, t+h/2-0.2, w-0.2, 0.4,
                 font_size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
add_rect(sl, 0, 0, 13.33, 2.0, C_ACCENT)
add_text(sl, "Privacy-Preserving Synthetic Student", 0.4, 0.15, 12.5, 0.85,
         font_size=32, bold=True, color=C_HIGHLIGHT, align=PP_ALIGN.CENTER)
add_text(sl, "Attendance Analysis and Prediction System", 0.4, 0.95, 12.5, 0.85,
         font_size=32, bold=True, color=C_HIGHLIGHT, align=PP_ALIGN.CENTER)
add_rect(sl, 1.5, 2.2, 10.33, 0.05, C_HIGHLIGHT)
add_bullet_box(sl, [
    "A Complete Data Science Project",
    "Data Generation  ·  EDA  ·  Regression  ·  Classification",
    "Dataset: 205 Students · 4,100 Records · Fully Synthetic",
    "Best Models: GradientBoostingRegressor  &  GradientBoostingClassifier",
    "Department: Computer Engineering & MCA  |  August 2026",
], 2.0, 2.45, 9.33, 3.0, font_size=19, body_color=C_LIGHT)
footer(sl, "Slide 1 / 14")
add_notes(sl, "Title slide. Emphasize upfront: all data is 100% synthetic. No real student info used.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCTION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Introduction", "Why student attendance prediction matters")
add_bullet_box(sl, [
    "•  Student attendance is a key indicator of academic engagement & performance",
    "•  Institutions enforce a mandatory 75% minimum for examination eligibility",
    "•  Identifying at-risk students early enables timely counseling intervention",
    "•  ML models can predict attendance status from academic & engagement features",
    "•  Real student data creates PII exposure & re-identification privacy risks",
    "•  Solution: fully synthetic data that mirrors realistic patterns ethically",
], 0.5, 1.5, 12.33, 5.0, font_size=18, body_color=C_LIGHT)
footer(sl, "Slide 2 / 14")
add_notes(sl, "Most colleges require 75% attendance. Early identification of defaulters is the motivation. "
              "Synthetic data solves the privacy challenge without sacrificing statistical realism.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Problem Statement", "Two machine learning tasks defined")
add_bullet_box(sl, [
    "•  Student absenteeism is often identified too late for effective intervention",
    "•  Faculty need a data-driven early-warning tool to flag at-risk students",
    "▶  Task 1 — Regression: Predict Attendance_Percentage (continuous 0–100%)",
    "▶  Task 2 — Classification: Predict Attendance_Status → Regular (≥75%) / Defaulter (<75%)",
    "•  Challenge: real attendance data contains PII — names, roll numbers, emails",
    "•  Solution: synthetic data preserves statistical realism without privacy cost",
], 0.5, 1.5, 12.33, 5.0, font_size=18,
   title_color=C_WARN, body_color=C_LIGHT)
footer(sl, "Slide 3 / 14")
add_notes(sl, "Two clear problem definitions. Regression for continuous prediction, "
              "classification for binary flag. Privacy constraint solved by synthetic data.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OBJECTIVES
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Project Objectives")
add_bullet_box(sl, [
    "1.  Generate a statistically realistic, fully synthetic 205-student dataset",
    "2.  Validate with 25 automated data quality & privacy checks — all must pass",
    "3.  Perform EDA with 11 professional visualizations & correlation analysis",
    "4.  Train & compare 4 regression models → predict Attendance_Percentage",
    "5.  Train & compare 4 classification models → predict Regular / Defaulter",
    "6.  Select best models using principled metric-based criteria (not arbitrary)",
], 0.5, 1.5, 12.33, 5.0, font_size=18, body_color=C_LIGHT)
footer(sl, "Slide 4 / 14")
add_notes(sl, "Six objectives spanning the complete data science lifecycle. "
              "Emphasize objective 6: model selection uses RMSE for regression, "
              "F1-score + Recall for classification.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DATASET GENERATION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Synthetic Dataset Generation", "data/student_attendance_205_students.csv")
add_bullet_box(sl, [
    "•  4,100 rows × 21 columns  |  205 students (STU0001–STU0205)",
    "•  CE cohort (STU0001–STU0060): Third Year, Fifth Semester — seed=42",
    "•  MCA cohort (STU0061–STU0205): Final Year, Third Semester — seed=242",
    "•  20 records/student: 5 subjects × 4 attendance periods",
    "•  Attendance rule: Regular ≥ 75% | Defaulter < 75%",
    "•  Validation: 25 checks, ALL 25 PASSED (validate_final_dataset.py)",
], 0.5, 1.5, 6.5, 5.0, font_size=16, body_color=C_LIGHT)
# stats box
add_rect(sl, 7.2, 1.55, 5.7, 4.8, RGBColor(0x0F, 0x3C, 0x78))
add_bullet_box(sl, [
    "Dataset Statistics",
    "Mean Attendance  :  68.67%",
    "Median           :  70.00%",
    "Std Deviation    :  20.21%",
    "Regular Records  :  1,834  (44.73%)",
    "Defaulter Records:  2,266  (55.27%)",
    "Missing Values   :  0",
    "Duplicate Rows   :  0",
], 7.3, 1.65, 5.5, 4.6, font_size=15,
   title_color=C_WARN, body_color=C_LIGHT)
footer(sl, "Slide 5 / 14")
add_notes(sl, "Dataset built in two phases. All values from the validated CSV. "
              "25 checks cover file existence, row counts, ID ranges, missing values, "
              "duplicate check, attendance calculation accuracy, and privacy checks.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PRIVACY & ETHICS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Data Privacy and Ethics", "Six privacy principles implemented")
add_rect(sl, 0.5, 1.5, 12.33, 0.6, RGBColor(0x8B, 0x00, 0x00))
add_text(sl,
    '"This project uses a synthetic student attendance dataset created for academic demonstration. '
    'It does not contain real student names, roll numbers, email IDs, or actual attendance '
    'records of identifiable students."',
    0.55, 1.55, 12.2, 0.5, font_size=12, bold=True, color=C_WARN, align=PP_ALIGN.CENTER)
add_bullet_box(sl, [
    "•  Data Minimization — no personal data collected or stored",
    "•  Full Anonymization — STU0001–STU0205; no real names or roll numbers",
    "•  No Re-identification — no mapping table created between IDs and real students",
    "•  Synthetic Labels — every file, notebook, and chart is clearly labeled SYNTHETIC",
    "•  Source Isolation — private_original_data/ never read by any project script",
    "•  No Discrimination — Defaulter label used for academic analysis only",
], 0.5, 2.25, 12.33, 4.5, font_size=17, body_color=C_LIGHT)
footer(sl, "Slide 6 / 14")
add_notes(sl, "Privacy is a design principle, not an afterthought. All six principles "
              "were implemented from day one. Emphasize the source isolation principle.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DATASET COLUMNS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Dataset Columns — 21 Features", "Feature categorization for modelling")
add_bullet_box(sl, [
    "▶  REGRESSION TARGET:  Attendance_Percentage  (continuous 0–100%)",
    "▶  CLASSIFICATION TARGET:  Attendance_Status  (Regular / Defaulter)",
], 0.5, 1.5, 12.33, 0.8, font_size=17, title_color=C_WARN, body_color=C_LIGHT)
add_bullet_box(sl, [
    "CATEGORICAL PREDICTORS (5 used in models):",
    "  Gender · Department · Year · Semester · Subject",
    "",
    "NUMERIC PREDICTORS (10 used in models):",
    "  Age · Previous_Attendance_Percentage · Assignment_Score · Internal_Marks",
    "  Study_Hours_Per_Week · Medical_Leave_Days · Travel_Distance_KM",
    "  Previous_Exam_Score · Late_Count · Online_Class_Attendance",
    "",
    "EXCLUDED — DATA LEAKAGE:  Classes_Attended · Total_Classes",
    "  (mathematically derived from Attendance_Percentage — including would be circular)",
], 0.5, 2.45, 12.33, 4.3, font_size=15, body_color=C_LIGHT)
footer(sl, "Slide 7 / 14")
add_notes(sl, "Classes_Attended is excluded because Attendance_Percentage = "
              "Classes_Attended / Total_Classes × 100. Including it would let the model "
              "reverse-compute the target — data leakage.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PREPROCESSING
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Data Preprocessing Pipeline", "Shared by regression & classification")
add_bullet_box(sl, [
    "•  Imputation  —  Numeric: median  |  Categorical: most_frequent",
    "•  Encoding   —  Categorical: OneHotEncoder (handle_unknown='ignore')",
    "•  Scaling    —  Numeric: StandardScaler (zero mean, unit variance)",
    "•  Split      —  80% Train / 20% Test  |  random_state=42  |  Stratified for classification",
    "•  Pipeline   —  sklearn Pipeline: Preprocessor → Estimator (single reusable object)",
], 0.5, 1.5, 12.33, 2.8, font_size=17, body_color=C_LIGHT)
add_rect(sl, 0.5, 4.4, 12.33, 2.6, RGBColor(0x0F, 0x3C, 0x78))
add_bullet_box(sl, [
    "Pipeline Architecture:",
    "  15 Predictors → ColumnTransformer",
    "       ├── Numeric Branch:  SimpleImputer(median) → StandardScaler",
    "       └── Categorical Branch:  SimpleImputer(most_frequent) → OneHotEncoder",
    "  → Regressor / Classifier → Predictions",
], 0.6, 4.5, 12.13, 2.4, font_size=14, title_color=C_WARN, body_color=C_LIGHT)
footer(sl, "Slide 8 / 14")
add_notes(sl, "The Pipeline object is critical: it ensures preprocessing is fitted "
              "only on training data, preventing leakage from test data statistics.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EDA
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Exploratory Data Analysis", "11 visualizations · Correlation analysis")
add_bullet_box(sl, [
    "•  Mean: 68.67%  |  Median: 70.00%  |  Std Dev: 20.21%",
    "•  Defaulter: 2,266 records (55.27%)  |  Regular: 1,834 records (44.73%)",
    "•  Subject variation < 1% — attendance is student-driven, not subject-driven",
    "•  Prior behavior dominates: Previous_Attendance_Pct (r=0.72), Internal_Marks (r=0.67)",
    "•  Logistical factors: Travel_Distance_KM (r=0.04) — no meaningful correlation",
], 0.5, 1.5, 6.8, 4.0, font_size=16, body_color=C_LIGHT)
try_add_image(sl, "correlation_heatmap.png", 7.2, 1.45, 5.8, 5.5)
footer(sl, "Slide 9 / 14")
add_notes(sl, "Key EDA insight: prior academic behavior predicts current attendance. "
              "Logistical factors (travel, medical leave, late count) have near-zero correlation. "
              "Correlation ≠ causation — especially important for synthetic data.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — REGRESSION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Regression Modelling", "Target: Attendance_Percentage (continuous)")
add_bullet_box(sl, [
    "4 Models: LinearRegression · DecisionTree · RandomForest · GradientBoosting",
    "Selection Criterion: Lowest RMSE + Highest R²",
    "",
    "  Model                      MAE        RMSE       R²",
    "  LinearRegression           9.3679     11.6108    0.6755",
    "  DecisionTreeRegressor      10.2138    12.8313    0.6037",
    "  RandomForestRegressor      9.1673     11.4357    0.6852",
    "  GradientBoostingRegressor  9.1786     11.3952    0.6874  ← BEST",
], 0.5, 1.5, 6.8, 5.2, font_size=14, title_color=C_WARN, body_color=C_LIGHT)
try_add_image(sl, "regression_actual_vs_predicted.png", 7.2, 1.45, 5.8, 5.5)
footer(sl, "Slide 10 / 14  |  Source: outputs/regression_model_results.csv")
add_notes(sl, "GradientBoosting achieves lowest RMSE (11.3952) and highest R² (0.6874). "
              "R² = 0.6874 means 68.74% of variance explained. Remaining ~31% comes from "
              "features with near-zero correlation (travel, medical, late count).")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLASSIFICATION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Classification Modelling", "Target: Attendance_Status — Regular / Defaulter")
add_bullet_box(sl, [
    "4 Models · Selection: Highest F1-Score, then Recall (NOT accuracy alone)",
    "Regular=0  |  Defaulter=1  |  Stratified 80/20 split",
    "",
    "  Model                       Acc     Prec    Recall  F1      AUC",
    "  LogisticRegression          0.8390  0.8512  0.8587  0.8549  0.9246",
    "  DecisionTreeClassifier      0.8268  0.8659  0.8124  0.8383  0.8841",
    "  RandomForestClassifier      0.8451  0.8655  0.8521  0.8587  0.9276",
    "  GradientBoostingClassifier  0.8463  0.8625  0.8587  0.8606  0.9189  ← BEST",
], 0.5, 1.5, 6.8, 5.2, font_size=13, title_color=C_WARN, body_color=C_LIGHT)
try_add_image(sl, "best_classifier_confusion_matrix.png", 7.2, 1.45, 5.8, 5.5)
footer(sl, "Slide 11 / 14  |  Source: outputs/classification_model_results.csv")
add_notes(sl, "GradientBoosting selected for highest F1 (0.8606). "
              "RandomForest has higher ROC-AUC (0.9276) but lower F1 — "
              "F1 is the primary criterion for early-warning systems. "
              "Recall=0.8587 means we correctly flag ~86% of all actual Defaulters.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — MODEL COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Final Model Comparison & Results", "Verified from actual CSV output files")
add_rect(sl, 0.5, 1.5, 12.33, 0.5, RGBColor(0x0F, 0x3C, 0x78))
add_text(sl, "All values sourced directly from outputs/regression_model_results.csv "
             "and outputs/classification_model_results.csv — no values estimated or fabricated.",
         0.6, 1.55, 12.1, 0.4, font_size=11, color=C_WARN, align=PP_ALIGN.CENTER)
add_bullet_box(sl, [
    "▶  BEST REGRESSION MODEL: GradientBoostingRegressor",
    "     MAE: 9.1786  |  MSE: 129.8510  |  RMSE: 11.3952  |  R²: 0.6874",
    "     Saved: models/best_regression_model.joblib",
    "",
    "▶  BEST CLASSIFICATION MODEL: GradientBoostingClassifier",
    "     Accuracy: 84.63%  |  Precision: 86.25%  |  Recall: 85.87%",
    "     F1-Score: 0.8606  |  ROC-AUC: 0.9189",
    "     Saved: models/best_classification_model.joblib",
    "",
    "•  Gradient Boosting outperforms on both tasks",
    "•  Linear/Logistic models are competitive — confirm value of interpretable baselines",
], 0.5, 2.1, 12.33, 4.8, font_size=16, title_color=C_WARN, body_color=C_LIGHT)
footer(sl, "Slide 12 / 14")
add_notes(sl, "Consolidate results. Emphasize: GradientBoosting wins both tasks. "
              "Margin over RandomForest is small — both are strong ensemble methods. "
              "Both model pipelines are saved and ready for reuse.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — LIMITATIONS & FUTURE SCOPE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Limitations and Future Scope")
add_bullet_box(sl, [
    "LIMITATIONS:",
    "•  Synthetic data — may not fully reflect real student behavior",
    "•  No causal analysis — correlations identified, not proven causes",
    "•  No temporal modeling — period totals, not week-by-week trends",
    "•  External factors absent: health, family, income not captured",
    "•  Retraining required before any real deployment",
], 0.5, 1.5, 6.3, 5.5, font_size=16, title_color=C_WARN, body_color=C_LIGHT)
add_bullet_box(sl, [
    "FUTURE SCOPE:",
    "•  Time-series LSTM for week-by-week deterioration detection",
    "•  SHAP explainability: WHY is a student flagged as Defaulter?",
    "•  Multi-class: Regular / At-Risk / Critical-Defaulter",
    "•  Streamlit dashboard for real-time teacher/HOD alerts",
    "•  Real data integration with proper institutional consent",
], 6.8, 1.5, 6.2, 5.5, font_size=16, title_color=C_HIGHLIGHT, body_color=C_LIGHT)
footer(sl, "Slide 13 / 14")
add_notes(sl, "The most significant limitation is synthetic data. But this is also "
              "the project's ethical strength. SHAP explainability is the most impactful "
              "near-term future scope for real institutional deployment.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Conclusion", "Phases 0–5 Complete")
add_bullet_box(sl, [
    "•  Synthetic dataset: 205 students · 4,100 records · 25/25 validation checks passed",
    "•  EDA: prior academic behavior dominates; logistical factors show near-zero correlation",
    "•  Best Regression:  GradientBoostingRegressor  |  RMSE: 11.3952  |  R²: 0.6874",
    "•  Best Classification:  GradientBoostingClassifier  |  F1: 0.8606  |  Recall: 0.8587",
    "•  Both model pipelines saved (.joblib) and ready for deployment or retraining",
    "•  Zero real student data — ethically compliant and safe for public submission",
], 0.5, 1.5, 12.33, 4.2, font_size=18, body_color=C_LIGHT)
add_rect(sl, 0.5, 5.9, 12.33, 0.9, RGBColor(0x0F, 0x3C, 0x78))
add_text(sl,
    '"This project demonstrates that a complete, production-quality attendance prediction system '
    'can be built ethically, using synthetic data that mirrors real-world statistical properties '
    '— without compromising any student\'s privacy."',
    0.6, 5.95, 12.1, 0.8, font_size=13, bold=True,
    color=C_WARN, align=PP_ALIGN.CENTER)
footer(sl, "Slide 14 / 14  |  Thank You")
add_notes(sl, "Conclude with two key takeaways: "
              "(1) Technical: GradientBoosting ensemble methods consistently outperform "
              "simpler models on mixed-feature attendance data. "
              "(2) Ethical: synthetic data is a viable, responsible substitute for real "
              "student records in academic ML projects. Thank the audience and invite questions.")

# ─── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
print(f"Slides: {len(prs.slides)}")
