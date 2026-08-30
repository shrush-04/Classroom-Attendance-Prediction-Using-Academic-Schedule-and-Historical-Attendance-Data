import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    # Use widescreen aspect ratio (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    c_dark_bg = RGBColor(17, 25, 40)       # #111928 (very dark blue/gray)
    c_white = RGBColor(255, 255, 255)
    c_gold = RGBColor(251, 191, 36)        # #fbbf24
    c_cyan = RGBColor(0, 242, 254)         # #00f2fe
    c_gray = RGBColor(143, 160, 181)       # #8fa0b5
    c_red = RGBColor(248, 113, 113)        # #f87171

    # Blank layout is layout index 6
    blank_layout = prs.slide_layouts[6]

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c_dark_bg

    def add_title(slide, text):
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Helvetica'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = c_cyan

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)
    
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Classroom Attendance Prediction"
    p1.font.name = 'Helvetica'
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = c_white
    
    p2 = tf.add_paragraph()
    p2.text = "Using Academic Schedule and Historical Attendance Data"
    p2.font.name = 'Helvetica'
    p2.font.size = Pt(24)
    p2.font.color.rgb = c_cyan
    p2.space_after = Pt(40)
    
    p3 = tf.add_paragraph()
    p3.text = "A Privacy-Preserving System for Lecture-Level Analysis\nMCA Final Year | Semester III Project\nDate: 2026-08-30"
    p3.font.name = 'Helvetica'
    p3.font.size = Pt(16)
    p3.font.color.rgb = c_gray

    # --- Slide 2: Scientific Status Overview ---
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)
    add_title(slide2, "Scientific Status & Core Findings")
    
    txBox = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Dataset Scale: Only 18 valid lecture observations are currently available."
    p.font.size = Pt(20)
    p.font.color.rgb = c_white
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Regression Status (Students_Present): Exploratory only. Model marginally beat the dummy baseline on a test split of only 4 observations. This is insufficient to establish generalization."
    p.font.size = Pt(20)
    p.font.color.rgb = c_gold
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Classification Status (Attendance_Band): Invalid for operational decisions. All models tied the dummy baseline accuracy of 0.50. The \"High\" attendance band (>75%) was never observed."
    p.font.size = Pt(20)
    p.font.color.rgb = c_red
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Operational Fallback: Historical-average baseline (38.75% attendance / ~31 students) is the primary recommendation for scheduling decisions."
    p.font.size = Pt(20)
    p.font.color.rgb = c_white

    # --- Slide 3: Problem Statement ---
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)
    add_title(slide3, "Problem Statement")
    
    txBox = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Privacy & Compliance Vulnerability: Storing individual student PII (names, roll numbers, emails) in prediction systems creates security risks and violates student privacy guidelines."
    p.font.size = Pt(20)
    p.font.color.rgb = c_white
    p.space_after = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• Proactive Timetable Management: Departments lack predictive, aggregate tools to understand slot-wise attendance drops before they happen."
    p.font.size = Pt(20)
    p.font.color.rgb = c_white
    p.space_after = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• The Lecture-Level Solution: Aggregate data at the lecture slot level. Planners only need to know when and why attendance drops, not who is absent."
    p.font.size = Pt(20)
    p.font.color.rgb = c_cyan

    # --- Slide 4: Data Collection & Privacy ---
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)
    add_title(slide4, "Data Collection & Privacy by Design")
    
    txBox = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Core Administrative Sources: Master Timetable (slots, classrooms), Faculty Registers (headcounts), Academic Calendar (holiday buffers), Continuous Evaluation Schedule (tests)."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Zero Student PII: No names, roll numbers, or college email IDs are collected, processed, or stored."
    p.font.size = Pt(18)
    p.font.color.rgb = c_cyan
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Faculty Anonymization: Instructor identifiers are mapped to codes (e.g. F_01, F_02) before modeling to protect teacher privacy."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Strict Source Integrity: All original files remain read-only inputs. Preprocessing outputs are written to dedicated pipelines."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white

    # --- Slide 5: Schema & Automated Validation ---
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)
    add_title(slide5, "Standardized Schema & Validation")
    
    txBox = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Schema Structure: 23 standardized columns mapping temporal features (Date, Day_of_Week), scheduling (Lecture_Number, Subject, Faculty_ID), and academic context (Test_Week, Assignment_Due, Holiday_Before_After)."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• 16 Automated Integrity Rules: Validated by validate_raw_data.py to enforce constraints like present count <= enrolled capacity, non-negative inputs, matching formats, and PII absence."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Key Patches Implemented:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = c_cyan
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  - Extended regex to support multi-faculty codes (e.g. F_01+F_13)\n  - Handled NaN values in Holiday_Before_After when read as empty strings\n  - Fixed PII check false-positives on 'Total_Enrolled_Students'"
    p.font.size = Pt(16)
    p.font.color.rgb = c_gray

    # --- Slide 6: Preprocessing & Leakage Prevention ---
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6)
    add_title(slide6, "Preprocessing & Leakage Prevention")
    
    txBox = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Leakage Elimination: Strict chronological sorting followed by shift operations (shift=1) on target-derived features to ensure future data is never leaked into training models."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Engineered Features (29 total):"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = c_cyan
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  - Previous_Lecture_Attendance_Percentage (Lag-1 value)\n  - Rolling_Average_Previous_3_Lectures (Moving trend)\n  - Subject_Historical_Average (Expanding window target encoding)\n  - Gap_Since_Previous_Lecture_Hours & Days_Since_Last_Holiday"
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Target Exclusion: Target columns (Students_Present and Attendance_Percentage) and Lecture_ID are strictly removed from features before training."
    p.font.size = Pt(18)
    p.font.color.rgb = c_gold

    # --- Slide 7: Exploratory Data Analysis (EDA) ---
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7)
    add_title(slide7, "Exploratory Data Analysis Summary")
    
    txBox = slide7.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Observations: 18 valid lecture sessions (2026-06-25 to 2026-08-07)"
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Core Stats: Mean Attendance: 38.75% (~31 students) | Range: 10.0% to 75.0% | Max Present: 60 students out of 80"
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Attendance Bands: Low (<50%): 12 lectures (67%) | Medium (50%-75%): 6 lectures (33%)"
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• ⚠️ High Attendance Class Not Represented: No lecture achieved attendance above 75%. The 'High' band (>75%) was never observed. Models cannot learn or predict this band."
    p.font.size = Pt(18)
    p.font.color.rgb = c_red
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Baseline Fallback: 38.75% average represents a realistic operational fallback for schedulers."
    p.font.size = Pt(18)
    p.font.color.rgb = c_cyan

    # --- Slide 8: Machine Learning Experiment Setup ---
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8)
    add_title(slide8, "Machine Learning Experiment Setup")
    
    txBox = slide8.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Chronological Splitting: First 80% (14 rows) for training, last 20% (4 rows) for testing. Random splits were avoided to prevent temporal leakage."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Tested Algorithms:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = c_cyan
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  - Regression: Linear Regression, Decision Tree, Random Forest, Gradient Boosting\n  - Classification: Logistic Regression, Decision Tree, Random Forest, SVM, k-NN"
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Evaluation Safeguard: Models must beat dummy baselines (mean-predictor for regression, most-frequent class for classification) on the test split to be considered valid."
    p.font.size = Pt(18)
    p.font.color.rgb = c_gold

    # --- Slide 9: Regression Results ---
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9)
    add_title(slide9, "Regression Results (Target: Students_Present)")
    
    txBox = slide9.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Dummy Baseline MAE: 14.5000 | Best Model (Random Forest) MAE: 14.0192"
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Random Forest Evaluation Metrics: RMSE: 15.1268 | R²: 0.1174 | MAPE: 43.05%"
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• ⚠️ Exploratory Classification: The regression model marginally beat the baseline (by 0.48 students) on a test split of only 4 observations. This result is exploratory only and cannot establish reliable generalization."
    p.font.size = Pt(18)
    p.font.color.rgb = c_gold
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Baseline Comparison: Random Forest explains only ~12% of the test variance (R²=0.1174), and predictions have an average relative error of 43%. All other regressors (LR, DT, GBR) were worse than baseline."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white

    # --- Slide 10: Classification Results ---
    slide10 = prs.slides.add_slide(blank_layout)
    set_background(slide10)
    add_title(slide10, "Classification Results (Target: Attendance_Band)")
    
    txBox = slide10.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Target Band Categories: Low (<50%), Medium (50%-75% inclusive), High (>75%)."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Dummy Classifier Accuracy: 0.5000 | Best Classifier (Logistic Regression) Accuracy: 0.5000"
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Weighted F1-score: 0.3333 | Accuracy: 0.5000 (Ties dummy baseline exactly)"
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• 🚫 Invalid for Operational Use: Since no model beat the baseline and the 'High' band was never observed, the classification model is invalid. Automated attendance-band decisions have been disabled."
    p.font.size = Pt(18)
    p.font.color.rgb = c_red
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Class Fallback: Predict historical most-frequent class (Low) as reference only."
    p.font.size = Pt(18)
    p.font.color.rgb = c_cyan

    # --- Slide 11: Streamlit Dashboard UI ---
    slide11 = prs.slides.add_slide(blank_layout)
    set_background(slide11)
    add_title(slide11, "Streamlit User Interface & Safety Warnings")
    
    txBox = slide11.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Historical Analysis (Tab 1): Visualizes trends, timeline graphs, and day-wise patterns from cleaned logs."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Predictive Inference (Tab 2): Renders future lecture predictions based on schedule input."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Prominent Scientific Warnings Rendered:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = c_cyan
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  - Insufficient history warning shows dataset size of 18\n  - Regression output is labeled as 'Exploratory / Limited Data'\n  - Attendance band prediction is flagged as unavailable and blocked\n  - Historical-average baseline (38.75%) is shown as the primary fallback prediction"
    p.font.size = Pt(16)
    p.font.color.rgb = c_gray

    # --- Slide 12: Project Limitations ---
    slide12 = prs.slides.add_slide(blank_layout)
    set_background(slide12)
    add_title(slide12, "Project Limitations & Data Constraints")
    
    txBox = slide12.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Tiny Sample Size: 18 total lecture sessions are insufficient for generalizable supervised learning."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• 4-Row Test Evaluation: The chronological 80/20 partition yields a 4-row test split, meaning evaluation metrics are statistically unstable."
    p.font.size = Pt(18)
    p.font.color.rgb = c_white
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Missing Class Representation: Zero observations in the 'High' attendance category, skewing the classifier to Low/Medium."
    p.font.size = Pt(18)
    p.font.color.rgb = c_red
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Data Integrity Constraint: No fake data generation, student PII fabrication, or interpolation was allowed, keeping the project scientifically honest."
    p.font.size = Pt(18)
    p.font.color.rgb = c_cyan

    # --- Slide 13: Conclusion ---
    slide13 = prs.slides.add_slide(blank_layout)
    set_background(slide13)
    add_title(slide13, "Conclusion")
    
    txBox = slide13.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = (
        "\"The available dataset contained only 18 valid lecture observations. The regression "
        "experiment produced a small improvement over the historical-average baseline, but the "
        "test set contained only four observations, so the result is exploratory and cannot "
        "establish reliable generalization. The classification model did not outperform the "
        "dummy baseline and should not be used for operational decisions. More physically "
        "verified lecture records are required before deploying a reliable predictive system.\""
    )
    p.font.name = 'Helvetica'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.font.italic = True
    p.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "— Official Project Conclusion Statement"
    p2.font.name = 'Helvetica'
    p2.font.size = Pt(16)
    p2.font.color.rgb = c_gray

    # Save presentation
    output_path = "classroom-attendance-schedule-project/presentation/project_presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == '__main__':
    create_presentation()
